import torch
import json
import re
import random
import os
import base64
import tempfile
from fastapi import FastAPI, HTTPException
from pydantic import BaseModel
from transformers import AutoTokenizer, AutoModelForCausalLM, BitsAndBytesConfig
from fastapi.middleware.cors import CORSMiddleware
import uvicorn

# Extractors
import easyocr
import docx
from pptx import Presentation
from PyPDF2 import PdfReader

# -----------------------------
# 1. MODEL CONFIGURATION
# -----------------------------
MODEL_NAME = "meta-llama/Llama-3.2-3B-Instruct"

print(f"Loading model {MODEL_NAME}...")

tokenizer = AutoTokenizer.from_pretrained(MODEL_NAME)
tokenizer.pad_token = tokenizer.eos_token

bnb_config = BitsAndBytesConfig(
    load_in_4bit=True,
    bnb_4bit_compute_dtype=torch.float16,
    bnb_4bit_use_double_quant=True,
    bnb_4bit_quant_type="nf4"
)

model = AutoModelForCausalLM.from_pretrained(
    MODEL_NAME,
    device_map="auto",
    quantization_config=bnb_config
)

# Initialize OCR
print("Loading EasyOCR Engine...")
ocr_reader = easyocr.Reader(['en'])

# -----------------------------
# 2. UNIVERSAL EXTRACTOR
# -----------------------------
def extract_text_from_file(file_path):
    if not os.path.exists(file_path):
        return ""
    
    ext = os.path.splitext(file_path)[1].lower()
    text = ""
    
    try:
        if ext in ['.jpg', '.jpeg', '.png']:
            result = ocr_reader.readtext(file_path, detail=0)
            text = " ".join(result)
        elif ext == '.docx':
            doc = docx.Document(file_path)
            text = "\n".join([p.text for p in doc.paragraphs])
        elif ext == '.pptx':
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"): text += shape.text + "\n"
        elif ext == '.pdf':
            reader = PdfReader(file_path)
            for page in reader.pages:
                text += page.extract_text() or ""
    except Exception as e:
        print(f"Extraction Error: {e}")
        return ""
        
    return text.strip()

def chunk_text(text, chunk_size=4000, overlap=500):
    if len(text) <= chunk_size: return [text]
    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunks.append(text[start:end])
        start += (chunk_size - overlap)
    return chunks

# -----------------------------
# 3. PROMPT LOGIC
# -----------------------------
def build_prompt(type, content, num_questions, difficulty):
    level_desc = {
        "Easy": "Introductory, direct questions. Test basic definitions and recall. The wrong options (distractors) should be obviously incorrect.",
        "Medium": "Require understanding concepts, not just definitions. Use plausible distractors that might trick someone with surface-level knowledge.",
        "Thinkable": "Analytical scenarios and logic. The differences between the correct answer and the distractors should be subtle.",
        "Hard": "Generate scenario-based or application-based questions. The differences between the correct answer and the distractors should be subtle."
    }.get(difficulty, "Balanced questions.")

    system_prompt = (
        "You are an Expert Educator and Quiz Master. "
        "Test conceptual understanding strictly based on the provided content. "
        "Include a valid 'explanation' for each answer."
    )
    
    instruction = f"""
Generate EXACTLY {num_questions} MCQs for: {content[:15000]}
Difficulty: {difficulty} ({level_desc})

FORMAT (Strict JSON):
{{
  "questions": [
    {{
      "questionText": "...",
      "options": ["...", "...", "...", "..."],
      "correctAnswer": "...",
      "explanation": "Brief explanation of why it is correct.",
      "points": 10,
      "type": "multiple-choice"
    }}
  ]
}}

RULES:
1. NO MATH: Use logic, not arithmetic.
2. EXPLANATIONS: You MUST include the 'explanation' field.
"""
    return f"<|begin_of_text|><|start_header_id|>system<|end_header_id|>\n\n{system_prompt}<|eot_id|><|start_header_id|>user<|end_header_id|>\n\n{instruction}<|eot_id|><|start_header_id|>assistant<|end_header_id|>\n\n"

# -----------------------------
# 4. FASTAPI SERVICE
# -----------------------------
app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

class GeneratorRequest(BaseModel):
    type: str # 'topic', 'pdf', 'docx', 'pptx', 'image'
    content: str # Can be raw text OR a local file path
    count: int = 5
    difficulty: str = "Medium"

@app.post("/generate")
async def generate_questions(req: GeneratorRequest):
    # Handle File Extraction if content is a path or base64 data
    source_text = req.content
    if req.type in ['pdf', 'docx', 'pptx', 'image']:
        if os.path.exists(req.content):
            source_text = extract_text_from_file(req.content)
        elif req.content.startswith("base64:"):
            try:
                b64_data = req.content[7:]
                file_data = base64.b64decode(b64_data)
                ext = req.type if req.type != 'image' else 'png'
                with tempfile.NamedTemporaryFile(delete=False, suffix=f".{ext}") as tmp:
                    tmp.write(file_data)
                    tmp_path = tmp.name
                
                source_text = extract_text_from_file(tmp_path)
                os.remove(tmp_path)
            except Exception as e:
                print(f"Base64 processing error: {e}")
                raise HTTPException(status_code=400, detail="Invalid base64 file content.")
    
    if not source_text or len(source_text) < 20:
        raise HTTPException(status_code=400, detail="Content too short or file unreadable.")

    # Chunking
    if len(source_text) > 4000:
        source_text = random.choice(chunk_text(source_text))

    max_tokens = int(req.count * 300)
    
    for i in range(2):
        prompt = build_prompt(req.type, source_text, req.count, req.difficulty)
        inputs = tokenizer(prompt, return_tensors="pt").to(model.device)

        with torch.no_grad():
            output = model.generate(
                input_ids=inputs["input_ids"],
                attention_mask=inputs["attention_mask"],
                max_new_tokens=max_tokens,
                temperature=0.4,
                do_sample=True,
                pad_token_id=tokenizer.eos_token_id
            )

        decoded = tokenizer.decode(output[0], skip_special_tokens=False)
        if "<|start_header_id|>assistant<|end_header_id|>" in decoded:
             decoded = decoded.split("<|start_header_id|>assistant<|end_header_id|>")[-1].replace("<|eot_id|>", "").strip()

        decoded = re.sub(r"```json", "", decoded).replace("```", "").strip()

        try:
            start = decoded.find('{')
            end = decoded.rfind('}') + 1
            data = json.loads(decoded[start:end])
            if "questions" in data: return data
        except:
            continue

    raise HTTPException(status_code=500, detail="AI generation failed.")

if __name__ == "__main__":
    uvicorn.run(app, host="0.0.0.0", port=8000)
